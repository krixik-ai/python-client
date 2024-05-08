import os
from pptx import Presentation


def is_valid(local_file_path: str) -> None:
    if local_file_path.split(".")[-1] != "pptx":
        raise ValueError(f"file is not a pptx file - {local_file_path}")

    try:
        # check that local_file_path represents pptx file
        reader = Presentation(local_file_path)

        # loop through slides, shapes, paragraphs, and runs to get all text
        for slide in reader.slides:
            if slide is not None:
                return None
    except FileNotFoundError:
        raise FileNotFoundError(f"The file '{local_file_path}' does not exist.")
    except Exception as e:
        raise ValueError(f"pptx extraction failed with exception {e}")


def is_size(
    *,
    local_file_path: str,
    minimum_word_count: int = 10,
    minimum_file_size: float = 0.000001,
    maximum_file_size: float = 100.000001,
) -> None:
    # proper size
    def compute_size(file_path: str):
        try:
            # Get the size of the file in bytes
            file_size_bytes = os.path.getsize(file_path)

            # Convert the size to megabytes (MB)
            file_size_mb = file_size_bytes / (1024 * 1024)
            return file_size_mb
        except Exception as e:
            raise ValueError(f"pptx size calculation failed with exception {e}")

    # compute word count
    def compute_word_count(file_path: str) -> int:
        try:
            # check that local_file_path represents pptx file
            reader = Presentation(local_file_path)

            # loop through slides, shapes, paragraphs, and runs to get all text
            word_count = 0
            for slide in reader.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run is None:
                                pass
                            else:
                                line = run.text
                                words = line.split()
                                word_count += len(words)
            return word_count
        except Exception as e:
            raise ValueError(f"pptx word count calculation failed with exception {e}")

    # check size of input text file
    try:
        # check that local_file_path represents a valid text file
        is_valid(local_file_path)

        # compute file size in megabytes
        file_size = compute_size(local_file_path)

        # check that file size in megabytes is greater than minimum_file_size and less than maximum_file_size
        if file_size < minimum_file_size or file_size > maximum_file_size:
            raise ValueError(
                f"input file size is {file_size} megabytes: either less than {minimum_file_size} megabytes (current minimum size allowable) or greater than {maximum_file_size} megabytes (current maximum size allowable) - {local_file_path}"
            )

        # compute word count
        file_word_count = compute_word_count(local_file_path)

        # check that word count is greater than minimum_word_count
        if file_word_count < minimum_word_count:
            raise ValueError(f"file word count is less than {minimum_word_count} words (current minimum word count allowable) - {local_file_path}")

    except ValueError as ve:
        raise ve

    except Exception as e:
        raise ValueError(f"pptx extraction failed with exception {e}")
