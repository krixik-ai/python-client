import os
from pptx import Presentation
from krixik.utilities.validators.system.base.utilities.decorators import (
    file_converters_input_check,
)
from krixik.utilities.validators.data.pptx import is_size as is_size_pptx
from krixik.utilities.validators.data.text import is_size as is_size_text
from krixik.utilities.utilities import vprint


def delete_file(*, file_path: str, exception: str):
    if file_path is not None:
        if os.path.exists(file_path):
            os.remove(file_path)
        raise ValueError(
            f"there seems to be a problem with your pptx - file conversion failed with exception {str(exception)} - please check your pptx file of the converter documentation - python-pptx - for further information: https://python-pptx.readthedocs.io/en/stable/"
        )


@file_converters_input_check
def extract(
    *,
    local_file_path: str,
    local_save_directory: str | None = None,
    verbose: bool = True,
) -> str | None:
    if local_file_path is None:
        raise NameError("local_file_path not defined")
    if os.path.exists(local_file_path) is False:
        raise FileNotFoundError(f"the file '{local_file_path}' does not exist.")
    if local_save_directory is None:
        raise NameError("local_save_directory is not defined")
    if os.path.exists(local_save_directory) is False:
        raise FileNotFoundError(f"the directory '{local_save_directory}' does not exist.")
    if not isinstance(verbose, bool):
        raise ValueError("verbose is not a boolean")

    convert_save_path = ""
    try:
        file_name = "krixik_converted_version_" + os.path.splitext(os.path.basename(local_file_path))[0]
        extension = ".txt"
        convert_save_path = os.path.join(local_save_directory, file_name + extension)

        # remove file if it already exists
        if os.path.exists(convert_save_path):
            os.remove(convert_save_path)

        # creating a pptx reader object
        prs = Presentation(local_file_path)

        # loop through slides, shapes, paragraphs, and runs to get all text
        num = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        line = run.text
                        if num == 0:
                            with open(convert_save_path, "a", encoding="utf-8") as new_file:
                                new_file.write(line)
                        if num > 0:
                            if os.path.exists(convert_save_path):
                                with open(convert_save_path, "a", encoding="utf-8") as new_file:
                                    new_file.write(line)
                            else:
                                raise FileNotFoundError(f"file {convert_save_path} does not exist")

                        num += 1
        # report success
        vprint(
            f"SUCCESS: File conversion complete python-pptx, result saved to: {convert_save_path}",
            verbose=verbose,
        )

        # report size check
        vprint(
            "INFO: Checking that file size falls within acceptable parameters...",
            verbose=verbose,
        )
        is_size_text(local_file_path=convert_save_path)
        vprint("INFO:...success!", verbose=verbose)

        return convert_save_path
    except ValueError as ve:
        vprint(
            "INFO:...failure!  The converted (text) file does not fall within acceptable size parameters.",
            verbose=verbose,
        )
        delete_file(file_path=convert_save_path, exception=str(ve))
    except Exception as e:
        delete_file(file_path=convert_save_path, exception=str(e))


@file_converters_input_check
def convert(
    *,
    local_file_path: str | None,
    local_save_directory: str | None = None,
    verbose: bool = True,
) -> str:
    """pptx converter - converts pptx file to txt using the pptx library

    Parameters
    ----------
    local_file_path : str | None
        path to local file to convert
    local_save_directory : str | None
        local directory to save converted file
    verbose : bool, optional
        by default True

    Returns
    -------
    str | None
        if conversion successful returns path to converted file
    """

    if local_file_path is None:
        raise NameError("local_file_path not defined")
    if os.path.exists(local_file_path) is False:
        raise FileNotFoundError(f"the file '{local_file_path}' does not exist.")

    if "." + local_file_path.split(".")[-1] == ".pptx":
        if local_save_directory is None:
            raise NameError("local_save_directory is not defined")

        vprint("INFO: Converting pptx to text...", verbose=verbose)

        # check size and validity of input pptx file
        is_size_pptx(local_file_path=local_file_path)

        # convert pptx to text
        new_local_file_path = extract(
            local_file_path=local_file_path,
            local_save_directory=local_save_directory,
            verbose=verbose,
        )

        return new_local_file_path
    return local_file_path
